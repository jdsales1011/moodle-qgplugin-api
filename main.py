from flask import Flask
from flask import request

import sys
import os
import json
import io
import base64

import openai
import PyPDF2
import docx
from pptx import Presentation

openai.api_key = os.getenv("OPENAI_API_KEY")


app = Flask(__name__)
@app.route("/qgplugin/api/", methods = ["POST"])
def get_questions():
    # Get JSON parameter
    json_data = request.get_data(as_text=True)

    # Decode JSON data into dictionary
    params = json.loads(json_data)

    number = int(params['number'])
    q_type = int(params['type'])
    files = params['files']

    file_content = ''

    for file in files:
        file_extension = file['file_name'].split('.')[-1]
        file_decoded = base64.b64decode(file['file_encoded'])
        print(file['file_name'])
        print(file_extension)

        file_content += "\n"

        if (file_extension == "txt"):
            file_content += file_decoded.decode('utf-8')

        elif (file_extension == "pdf"):
            # Convert contents to bytes (BytesIO obj)
            file_bytes = io.BytesIO(file_decoded)
            pdf_reader = PyPDF2.PdfReader(file_bytes)
            for page in range(len(pdf_reader.pages)):
                file_content += pdf_reader.pages[page].extract_text()

            # Close BytesIO obj
            file_bytes.close()

        elif (file_extension == "docx" or file_extension == "docs" or file_extension == "doc"):
            file_bytes = io.BytesIO(file_decoded)
            doc = docx.Document(file_bytes)
            file_content += "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            # Close BytesIO obj
            file_bytes.close()

        elif (file_extension == "ppt" or file_extension == "pptx"):
            # Convert contents to bytes (BytesIO obj)
            file_bytes = io.BytesIO(file_decoded)

            prs = Presentation(file_bytes)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        file_content += shape.text + "\n"

            # Close BytesIO obj
            file_bytes.close()

    prompt_creators = {
        1: prompt_creator1(file_content, number),
        2: prompt_creator2(file_content, number),
        3: prompt_creator3(file_content, number),
    }
    prompt = prompt_creators[q_type]
    print(prompt)

    # PREDICT FUNCTION
    return predict_questions(prompt, q_type, number)


def prompt_creator1(content, num=1): #identification
    return f"Generate {num} Questions and their short answers as a list from the following text, {content} using the format 'Q1.' and 'Answer:', Answer must be a direct, concise, short answer. \n\nQuestions and Answers:"

def prompt_creator2(content, num=1): #true or false
    return f"Generate {num} true or false questions and their answers as a list from the following text, {content} using the format 'Q1.' and 'Answer:'. Phrase statements in declarative form. Answers should be either True or False.\n\Statements and Answers:"

def prompt_creator3(content, num=1): #multiple choice
    return f"Create a list of {num} multiple choice questions and their answers from the following text, {content} using the format 'Q1.','Choices: []' and 'Answer:'\n\nQuestions and Answers:"


# AI PREDICT FUNCTION
def predict_questions(prompt, q_type, number):
    result = []

    try:
        openai.api_key = os.getenv("OPENAI_API_KEY")    
        response = openai.Completion.create(
            model="text-davinci-003",
            prompt=prompt,
            temperature=0.7,
            max_tokens=256,
            top_p=1,
            frequency_penalty=0,
            presence_penalty=0
        )
        new = response['choices'][0]['text']

        result = new.split('\n')

    except openai.OpenAIError as e:
        # Handle OpenAI API specific errors
        return generate_error(500, f"OpenAI API error: {str(e)}")

    except Exception as e:
        # Handle other general exceptions
        return generate_error(500, f"An error occurred while handling OpenAI request: : {str(e)}")


    try:
        ques_bank = []

        result = [x for x in result if x.strip()]       # TO REMOVE EMPTY NEXT LINES

        length = len(result)

        if q_type == 3:             # MULTIPLE CHOICE
            for i in range(0, length, 3):
                item = {
                    "id": int(i/3)+1,
                    "question" : result[i].split(".", 1)[1].strip(),
                    "choices" : result[i+1].split("Choices:", 1)[1].strip(),
                    "answer" : result[i+2].split("Answer:", 1)[1].split(".", 1)[0].strip()
                    }
                ques_bank.append(item)

        else: 
            for i in range(0, length, 2):
                item = {
                    "id": int(i/2)+1,
                    "question" : result[i].split(".", 1)[1].strip(),
                    "answer" : result[i+1].split("Answer:", 1)[1].strip()
                    }
                ques_bank.append(item)

        ques_bank_json = {
            "status" : "success",
            "questions": ques_bank,
        }

        # Convert dict to JSON
        ques_bank_json = json.dumps(ques_bank_json)

        print("JSON RESULT: \n", ques_bank_json)

        # Return a JSON 
        return ques_bank_json
    
    except Exception as e:
        # Handle other general exceptions
        json_response = generate_error(500, f"An error occured during handling of response: {str(e)}")
        print("RESPONSE: \n", json_response)
        return json_response


def generate_error(code, message):
    error_json = {
        "status" : "error",
        "error" : {
            "code" : code,
            "message" : message,
        }
    }

    # Convert dict to JSON
    error_json = json.dumps(error_json)

    return error_json


@app.route("/qgplugin/api/test/", methods = ["GET"])
def test_print():
    print("INSIDE API")
    return "INSIDE API !!!!"

if __name__ == '__main__':
    app.run(debug=True, port=os.getenv("PORT", default=5000))
